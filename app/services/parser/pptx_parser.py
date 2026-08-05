"""PPTX 바이너리에서 텍스트를 추출한다.

슬라이드별로 도형(shape)의 텍스트를 순서대로 이어붙여 하나의 문자열로
반환한다. 레이아웃/이미지/애니메이션 보존은 다루지 않는다(YAGNI, 후속 단계).
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from app.core.errors import ParserError


def extract_text(data: bytes) -> str:
    """PPTX 바이트에서 슬라이드별 도형 텍스트를 순서대로 이어붙여 반환한다.

    Raises:
        ParserError: 손상된 파일이거나 PPTX 형식이 아닌 경우.
    """
    try:
        presentation = Presentation(BytesIO(data))
        slides_text: list[str] = []
        for slide in presentation.slides:
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text
            ]
            if texts:
                slides_text.append("\n".join(texts))
    except Exception as exc:
        raise ParserError(f"failed to parse PPTX: {exc}") from exc
    return "\n\n".join(slides_text)
