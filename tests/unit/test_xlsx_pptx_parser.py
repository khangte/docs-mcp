"""XLSX/PPTX 텍스트 추출 단위 테스트 (기존 pdf_parser/docx_parser 와 동일 패턴)."""

from __future__ import annotations

from io import BytesIO

import pytest
from openpyxl import Workbook
from pptx import Presentation

from app.core.errors import ParserError
from app.services.parser import pptx_parser, xlsx_parser


def _make_xlsx_bytes(rows: list[list[object]]) -> bytes:
    """지정한 행들을 담은 XLSX 바이트를 만든다."""
    workbook = Workbook()
    sheet = workbook.active
    for row in rows:
        sheet.append(row)
    buf = BytesIO()
    workbook.save(buf)
    return buf.getvalue()


def _make_pptx_bytes(slide_texts: list[str]) -> bytes:
    """지정한 텍스트를 각각 한 슬라이드의 제목으로 담은 PPTX 바이트를 만든다."""
    presentation = Presentation()
    layout = presentation.slide_layouts[0]
    for text in slide_texts:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = text
    buf = BytesIO()
    presentation.save(buf)
    return buf.getvalue()


# --- xlsx_parser.extract_text -------------------------------------------------


def test_xlsx_extract_text_joins_cell_values() -> None:
    """시트의 셀 텍스트가 추출된다."""
    data = _make_xlsx_bytes([["이름", "역할"], ["강민혁", "백엔드"]])

    text = xlsx_parser.extract_text(data)

    assert "강민혁" in text
    assert "백엔드" in text


def test_xlsx_extract_text_preserves_row_order() -> None:
    """행 순서대로 텍스트가 이어붙여진다."""
    data = _make_xlsx_bytes([["첫 행"], ["둘째 행"]])

    text = xlsx_parser.extract_text(data)

    assert text.index("첫 행") < text.index("둘째 행")


def test_xlsx_extract_text_raises_parser_error_on_corrupt_bytes() -> None:
    """손상된 XLSX 바이트는 ParserError 다."""
    with pytest.raises(ParserError):
        xlsx_parser.extract_text(b"not a real xlsx")


# --- pptx_parser.extract_text -------------------------------------------------


def test_pptx_extract_text_joins_slide_texts() -> None:
    """슬라이드 도형 텍스트가 추출된다."""
    data = _make_pptx_bytes(["첫 슬라이드 제목", "둘째 슬라이드 제목"])

    text = pptx_parser.extract_text(data)

    assert "첫 슬라이드 제목" in text
    assert "둘째 슬라이드 제목" in text


def test_pptx_extract_text_preserves_slide_order() -> None:
    """슬라이드 순서대로 텍스트가 이어붙여진다."""
    data = _make_pptx_bytes(["앞 슬라이드", "뒤 슬라이드"])

    text = pptx_parser.extract_text(data)

    assert text.index("앞 슬라이드") < text.index("뒤 슬라이드")


def test_pptx_extract_text_raises_parser_error_on_corrupt_bytes() -> None:
    """손상된 PPTX 바이트는 ParserError 다."""
    with pytest.raises(ParserError):
        pptx_parser.extract_text(b"not a real pptx")
